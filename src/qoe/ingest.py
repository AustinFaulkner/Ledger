"""Load and chunk a data room into retrievable text chunks.

Supports PDF, Excel (.xlsx + legacy .xls), CSV/TSV, Word (.docx), PowerPoint (.pptx),
RTF, JSON, Markdown and plain text. A corrupt file or a missing optional dependency
fails soft (0 chunks) rather than aborting an upload. Keeping ingestion separate means
new file types are added here without touching the analysis logic.
"""

from __future__ import annotations

import csv
import io
from dataclasses import dataclass, field
from pathlib import Path

SUPPORTED = {".pdf", ".xlsx", ".xls", ".csv", ".tsv", ".docx", ".rtf", ".pptx", ".json", ".md", ".txt"}


@dataclass
class Chunk:
    """A retrievable unit of text plus a citation back to its origin."""
    text: str
    source: str           # filename
    locator: str          # page number, sheet!cell-range, or paragraph index
    meta: dict = field(default_factory=dict)


def load_dir(root: str | Path) -> list[Chunk]:
    """Walk a directory and chunk every supported document under it."""
    root = Path(root)
    chunks: list[Chunk] = []
    for path in sorted(root.rglob("*")):
        if path.suffix.lower() in SUPPORTED:
            chunks.extend(load_file(path))
    return chunks


def load_file(path: Path) -> list[Chunk]:
    """Chunk a single file by type. Any failure — corrupt file, missing optional
    dependency, or unsupported type — returns an empty list rather than raising, so one
    bad file never aborts a data-room upload."""
    handler = {
        ".pdf": _chunk_pdf,
        ".xlsx": _chunk_xlsx,
        ".xls": _chunk_xls,
        ".csv": _chunk_delimited,
        ".tsv": _chunk_delimited,
        ".docx": _chunk_docx,
        ".rtf": _chunk_rtf,
        ".pptx": _chunk_pptx,
        ".json": _chunk_json,
        ".md": _chunk_text,
        ".txt": _chunk_text,
    }.get(path.suffix.lower())
    if not handler:
        return []
    try:
        return handler(path)
    except Exception:
        return []  # corrupt / unreadable / optional dependency absent → soft-fail


def _text_to_chunks(text: str, source: str) -> list[Chunk]:
    """Split prose into paragraph chunks (blank-line separated)."""
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paras:
        return [Chunk(text=text, source=source, locator="full")] if text.strip() else []
    return [Chunk(text=p, source=source, locator=f"para:{i}") for i, p in enumerate(paras)]


def _chunk_text(path: Path) -> list[Chunk]:
    return _text_to_chunks(path.read_text(encoding="utf-8", errors="ignore"), path.name)


def _chunk_pdf(path: Path) -> list[Chunk]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    out: list[Chunk] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            # TODO: split long pages into ~500-token windows for tighter citations.
            out.append(Chunk(text=text, source=path.name, locator=f"p.{i}"))
    return out


def _chunk_xlsx(path: Path) -> list[Chunk]:
    """One chunk per worksheet, rendered as CSV via openpyxl (modern .xlsx)."""
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    try:
        out: list[Chunk] = []
        for ws in wb.worksheets:
            buf = io.StringIO()
            writer = csv.writer(buf)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(["" if c is None else c for c in row])
            text = buf.getvalue().strip()
            if text:
                out.append(Chunk(text=text, source=path.name, locator=f"sheet:{ws.title}"))
        return out
    finally:
        wb.close()


def _chunk_xls(path: Path) -> list[Chunk]:
    """One chunk per worksheet for legacy .xls via xlrd (openpyxl cannot read .xls)."""
    import xlrd

    book = xlrd.open_workbook(str(path))
    out: list[Chunk] = []
    for sheet in book.sheets():
        buf = io.StringIO()
        writer = csv.writer(buf)
        for r in range(sheet.nrows):
            writer.writerow([sheet.cell_value(r, c) for c in range(sheet.ncols)])
        text = buf.getvalue().strip()
        if text:
            out.append(Chunk(text=text, source=path.name, locator=f"sheet:{sheet.name}"))
    return out


def _chunk_delimited(path: Path) -> list[Chunk]:
    # CSV/TSV are already the text we want — read verbatim (faithful to source rows).
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    return [Chunk(text=text, source=path.name, locator="table")] if text else []


def _chunk_docx(path: Path) -> list[Chunk]:
    from docx import Document

    doc = Document(str(path))
    out: list[Chunk] = []
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            out.append(Chunk(text=para.text, source=path.name, locator=f"para:{i}"))
    return out


def _chunk_rtf(path: Path) -> list[Chunk]:
    from striprtf.striprtf import rtf_to_text

    text = rtf_to_text(path.read_text(encoding="utf-8", errors="ignore"))
    return _text_to_chunks(text, path.name)


def _chunk_pptx(path: Path) -> list[Chunk]:
    """One chunk per slide — shape text + table rows."""
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[Chunk] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        lines.append(", ".join(cells))
        text = "\n".join(lines).strip()
        if text:
            out.append(Chunk(text=text, source=path.name, locator=f"slide:{i}"))
    return out


def _chunk_json(path: Path) -> list[Chunk]:
    import json as _json

    raw = path.read_text(encoding="utf-8", errors="ignore")
    try:
        text = _json.dumps(_json.loads(raw), indent=2, ensure_ascii=False)
    except Exception:
        text = raw
    return _text_to_chunks(text, path.name)
